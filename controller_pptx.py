import openai
from pptx import Presentation
import json
import os
import hashlib

import logging
logging.basicConfig(level=logging.INFO, filename="app.log")

cache_dir = "cache"

prompt_ppt = """
For every scenario I provide, I need at detailed response in a json format. Be critical and provide accurate facts and figures.

Here's a sample json body I want:

{
  "title": "Main Title",
  "subtitle": "Sub Heading"
  "slides": [
    {
      "id": 1,
      "title": "Title of Slide 1",
      "content": [
        {
          "type": "text",
          "value": "Some text on slide 1"
        },
        {
          "type": "bullet_points",
          "value": ["Bullet point 1", "Bullet point 2", "Bullet point 3"]
        }
      ]
    }
  ]
}
Nothing outside of this format should be returned.

There should be only one slide title, one subheading, and multiple slides. The first slide should always be the Agenda slide.
The second slide should contain a brief history of the topic and an introduction. The last slide should be the conclusion slide.
Generate slides with a mix of both text and bullet points.
Each slide should not exceed more than 250 words.. 
Only return the JSON object. Truncate anything after the outermost closing parentheses.
"""

def get_response(prompt, user, key):
    try: 
        logging.info("Called response from pptx")
        openai.api_key = key
        resp = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": user},
          ],
        max_tokens=3500,
        temperature=0
        )
        logging.info("Showing response from pptx")
        logging.info(resp)
    except Exception as e:
        logging.error(str(e))
        logging.info("API error")
        pass
    
    # turbo ChatCompletion resp
    return resp.to_dict()["choices"][0]["message"]["content"].replace("\n", "").strip()


def create_pptx_from_json(json_text, user):
    # Load the JSON
    data = json.loads(json_text)

    # Initialize the PowerPoint presentation
    prs = Presentation("./init_template.pptx")

    # Add the title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = data['title']
    subtitle = slide.placeholders[1]
    subtitle.text = data['subtitle']

    # Add the remaining slides
    for slide_data in data['slides']:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_data['title']
        for content_data in slide_data['content']:
            if content_data['type'] == 'text':
                tf = slide.shapes.placeholders[1].text_frame
                p = tf.add_paragraph()
                p.text = content_data['value']
                p.level = 0
            elif content_data['type'] == 'bullet_points':
                tf = slide.shapes.placeholders[1].text_frame
                for point in content_data['value']:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 1

    # Generate a unique filename based on user information
    user_hash = hashlib.md5(user.encode()).hexdigest()
    logging.info("Saving with user: " + user)
    logging.info("Saving with hash: " + user_hash)
    filename = f"presentation_{user_hash}.pptx"
    cache_file = os.path.join(cache_dir, filename)
    # Save the PowerPoint presentation
    prs.save(cache_file)


def generate_pptx(prompt, user, key, filename):
    
    logging.info(user)
    # Get the response from the GPT-3 model
    logging.info("Generating PowerPoint presentation...")
    resp = get_response(prompt_ppt, user, key)
    logging.info("Done!")
    logging.info(resp)
    
    # Create the PowerPoint presentation
    logging.info("Creating PowerPoint presentation...LOL")
    create_pptx_from_json(resp, user)
    # logging.info("PowerPoint presentation saved as "+filename+".pptx")